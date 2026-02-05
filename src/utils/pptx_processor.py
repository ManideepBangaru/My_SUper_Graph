"""
PPTX Processor - Extracts text and images from PowerPoint files.

Processes PPTX files into chunks suitable for LLM context, with images uploaded to S3.
"""

import asyncio
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Support both LangGraph Studio and FastAPI server imports
try:
    from utils.s3_operations import S3Operations
    from utils.pdf_processor import DocumentChunk, PageData
except ImportError:
    from src.utils.s3_operations import S3Operations
    from src.utils.pdf_processor import DocumentChunk, PageData


class PPTXProcessor:
    """
    Processes PPTX files by extracting text and images from slides.
    
    Text is chunked using RecursiveCharacterTextSplitter for optimal LLM context.
    Images are uploaded to S3 and referenced by key.
    """
    
    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        s3_ops: Optional[S3Operations] = None,
    ):
        """
        Initialize the PPTX processor.
        
        Args:
            chunk_size: Maximum size of each text chunk
            chunk_overlap: Overlap between chunks for context continuity
            s3_ops: S3Operations instance for image uploads
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.s3_ops = s3_ops
        
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
    
    def _extract_text_from_shape(self, shape) -> str:
        """
        Extract text from a shape, handling various shape types.
        
        Args:
            shape: A pptx shape object
            
        Returns:
            Extracted text from the shape
        """
        text_parts = []
        
        # Handle shapes with text frames
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = "".join(run.text for run in paragraph.runs)
                if paragraph_text.strip():
                    text_parts.append(paragraph_text)
        
        # Handle tables
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    text_parts.append(" | ".join(row_texts))
        
        return "\n".join(text_parts)
    
    def _extract_images_from_shape(self, shape) -> list[tuple[bytes, str]]:
        """
        Extract images from a shape.
        
        Args:
            shape: A pptx shape object
            
        Returns:
            List of (image_bytes, extension) tuples
        """
        images = []
        
        # Check if shape is a picture
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_bytes = image.blob
                # Get extension from content type (e.g., "image/png" -> "png")
                ext = image.ext
                images.append((image_bytes, ext))
            except Exception:
                # Skip images that can't be extracted
                pass
        
        return images
    
    def _extract_slides(self, pptx_bytes: bytes) -> list[PageData]:
        """
        Extract text and images from each slide of a PPTX.
        
        Args:
            pptx_bytes: Raw PPTX file content
            
        Returns:
            List of PageData objects, one per slide
        """
        slides_data = []
        
        # Open PPTX from bytes
        prs = Presentation(BytesIO(pptx_bytes))
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text_parts = []
            slide_images = []
            
            # Process each shape in the slide
            for shape in slide.shapes:
                # Extract text
                text = self._extract_text_from_shape(shape)
                if text:
                    slide_text_parts.append(text)
                
                # Extract images
                images = self._extract_images_from_shape(shape)
                slide_images.extend(images)
                
                # Handle grouped shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for sub_shape in shape.shapes:
                            text = self._extract_text_from_shape(sub_shape)
                            if text:
                                slide_text_parts.append(text)
                            images = self._extract_images_from_shape(sub_shape)
                            slide_images.extend(images)
                    except Exception:
                        pass
            
            # Combine all text from the slide
            slide_text = "\n\n".join(slide_text_parts)
            
            slides_data.append(PageData(
                page_num=slide_num,
                text=slide_text,
                images=slide_images,
            ))
        
        return slides_data
    
    async def _upload_images(
        self,
        images: list[tuple[bytes, str]],
        slide_num: int,
        user_id: str,
        thread_id: str,
        filename: str,
    ) -> list[str]:
        """
        Upload extracted images to S3.
        
        Args:
            images: List of (image_bytes, extension) tuples
            slide_num: Slide number for naming
            user_id: User identifier
            thread_id: Thread identifier
            filename: Original PPTX filename for organizing images
            
        Returns:
            List of S3 keys for uploaded images
        """
        if not self.s3_ops or not images:
            return []
        
        image_keys = []
        base_name = filename.rsplit(".", 1)[0]
        
        for img_idx, (image_bytes, ext) in enumerate(images):
            image_filename = f"{base_name}_slide{slide_num}_img{img_idx}.{ext}"
            
            try:
                result = await self.s3_ops.upload_file(
                    file_data=image_bytes,
                    filename=image_filename,
                    user_id=user_id,
                    thread_id=thread_id,
                )
                image_keys.append(result["key"])
                print(f"Uploaded image: {result['key']}")
            except Exception as e:
                # Log error but continue with other images
                print(f"Failed to upload image {image_filename}: {e}")
                continue
        
        return image_keys
    
    def _chunk_text(self, text: str) -> list[str]:
        """
        Split text into chunks for LLM context.
        
        Args:
            text: Full slide text
            
        Returns:
            List of text chunks
        """
        if not text or not text.strip():
            return []
        
        chunks = self.text_splitter.split_text(text)
        return chunks
    
    async def process_pptx(
        self,
        pptx_bytes: bytes,
        user_id: str,
        thread_id: str,
        filename: str,
    ) -> list[DocumentChunk]:
        """
        Process a PPTX file, extracting text chunks and uploading images.
        
        Args:
            pptx_bytes: Raw PPTX file content
            user_id: User identifier for S3 organization
            thread_id: Thread identifier for S3 organization
            filename: Original filename of the PPTX
            
        Returns:
            List of DocumentChunk objects ready for database storage
        """
        # Extract slides in a thread to avoid blocking
        slides_data = await asyncio.to_thread(self._extract_slides, pptx_bytes)
        
        all_chunks = []
        
        for slide_data in slides_data:
            # Upload images for this slide
            print(f"[PPTX Processor] Slide {slide_data.page_num}: {len(slide_data.images)} images to upload")
            image_keys = await self._upload_images(
                images=slide_data.images,
                slide_num=slide_data.page_num,
                user_id=user_id,
                thread_id=thread_id,
                filename=filename,
            )
            print(f"[PPTX Processor] Slide {slide_data.page_num}: uploaded image_keys = {image_keys}")
            
            # Chunk the text
            text_chunks = self._chunk_text(slide_data.text)
            
            if text_chunks:
                # Create DocumentChunk for each text chunk
                for chunk_idx, chunk_text in enumerate(text_chunks):
                    # Only attach image keys to the first chunk of each slide
                    chunk_image_keys = image_keys if chunk_idx == 0 else []
                    
                    all_chunks.append(DocumentChunk(
                        page_num=slide_data.page_num,
                        chunk_index=chunk_idx,
                        content=chunk_text,
                        image_keys=chunk_image_keys,
                    ))
            elif image_keys:
                # Slide has images but no text - create a placeholder chunk
                all_chunks.append(DocumentChunk(
                    page_num=slide_data.page_num,
                    chunk_index=0,
                    content=f"[Slide {slide_data.page_num + 1}: Contains {len(image_keys)} image(s)]",
                    image_keys=image_keys,
                ))
        
        return all_chunks


async def process_uploaded_pptx_file(
    s3_key: str,
    user_id: str,
    thread_id: str,
    filename: str,
) -> int:
    """
    Background task to process an uploaded PPTX file.
    
    Downloads from S3, processes, and stores chunks in the database.
    
    Args:
        s3_key: S3 key of the uploaded file
        user_id: User identifier
        thread_id: Thread identifier
        filename: Original filename
        
    Returns:
        Number of chunks created
    """
    print(f"[PPTX Processor] Starting processing for {filename}")
    
    # Import here to avoid circular imports
    try:
        from api.database import save_document_chunks
    except ImportError:
        from src.api.database import save_document_chunks
    
    try:
        s3_ops = S3Operations()
        
        # Download the PPTX from S3
        print(f"[PPTX Processor] Downloading {filename} from S3...")
        pptx_bytes = await s3_ops.download_file(filename, user_id, thread_id)
        print(f"[PPTX Processor] Downloaded {len(pptx_bytes)} bytes")
        
        # Process the PPTX
        processor = PPTXProcessor(s3_ops=s3_ops)
        chunks = await processor.process_pptx(
            pptx_bytes=pptx_bytes,
            user_id=user_id,
            thread_id=thread_id,
            filename=filename,
        )
        print(f"[PPTX Processor] Extracted {len(chunks)} chunks")
        
        # Save to database
        chunk_dicts = [chunk.to_dict() for chunk in chunks]
        saved_count = await save_document_chunks(
            thread_id=thread_id,
            user_id=user_id,
            filename=filename,
            chunks=chunk_dicts,
        )
        
        print(f"[PPTX Processor] Saved {saved_count} chunks to database for {filename}")
        return saved_count
    except Exception as e:
        print(f"[PPTX Processor] Error processing {filename}: {e}")
        import traceback
        traceback.print_exc()
        raise
